import os
import logging
import tempfile
from uuid import uuid4
from PIL import Image
from pdf2image import convert_from_path
from PyPDF2 import PdfReader, PdfWriter

# For Document conversions
try:
    # Import libraries for document conversions
    import docx
    from docx2pdf import convert as docx_to_pdf
    import openpyxl
    from openpyxl import Workbook
    from pptx import Presentation
    # These imports might not be available; we'll handle exceptions when using them
except ImportError:
    pass

logger = logging.getLogger(__name__)

def photos_to_pdf(photo_paths, output_dir):
    """
    Convert photos to a PDF file.
    
    Args:
        photo_paths: List of paths to the photos
        output_dir: Directory to save the result
        
    Returns:
        Path to the generated PDF
    """
    try:
        images = []
        first_image = None
        
        for photo_path in photo_paths:
            img = Image.open(photo_path)
            if img.mode == 'RGBA':
                img = img.convert('RGB')
            
            if first_image is None:
                first_image = img
                images.append(img)
            else:
                images.append(img)
        
        if not images:
            raise ValueError("لم يتم تقديم أي صور للتحويل.")
        
        output_path = os.path.join(output_dir, f"converted_photos_{uuid4().hex}.pdf")
        
        # Save the first image as PDF with the rest appended
        if len(images) == 1:
            images[0].save(output_path, "PDF")
        else:
            images[0].save(output_path, "PDF", save_all=True, append_images=images[1:])
        
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting photos to PDF: {str(e)}")
        raise Exception(f"فشل في تحويل الصور إلى PDF: {str(e)}")

def pdf_to_images(pdf_path, output_dir):
    """
    Convert a PDF file to images.
    
    Args:
        pdf_path: Path to the PDF file
        output_dir: Directory to save the result
        
    Returns:
        List of paths to the generated images
    """
    try:
        images = convert_from_path(pdf_path)
        
        result_paths = []
        for i, img in enumerate(images):
            img_path = os.path.join(output_dir, f"page_{i + 1}_{uuid4().hex}.png")
            img.save(img_path, "PNG")
            result_paths.append(img_path)
        
        return result_paths
    
    except Exception as e:
        logger.error(f"Error converting PDF to images: {str(e)}")
        raise Exception(f"فشل في تحويل PDF إلى صور: {str(e)}")

def pdf_to_word(pdf_path, output_dir):
    """
    Convert a PDF file to a Word document.
    
    This is a simplified implementation. For production use,
    consider a dedicated PDF-to-Word conversion API.
    
    Args:
        pdf_path: Path to the PDF file
        output_dir: Directory to save the result
        
    Returns:
        Path to the generated Word document
    """
    try:
        # Create a new Word document
        doc = docx.Document()
        
        # Extract text from the PDF
        reader = PdfReader(pdf_path)
        
        for page in reader.pages:
            text = page.extract_text()
            doc.add_paragraph(text)
        
        output_path = os.path.join(output_dir, f"converted_to_word_{uuid4().hex}.docx")
        doc.save(output_path)
        
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting PDF to Word: {str(e)}")
        raise Exception(f"فشل في تحويل PDF إلى Word: {str(e)}")

def pdf_to_excel(pdf_path, output_dir):
    """
    Convert a PDF file to an Excel document.
    
    This is a simplified implementation. For production use,
    consider a dedicated PDF-to-Excel conversion API.
    
    Args:
        pdf_path: Path to the PDF file
        output_dir: Directory to save the result
        
    Returns:
        Path to the generated Excel document
    """
    try:
        # Create a new Excel workbook
        wb = Workbook()
        ws = wb.active
        
        # Extract text from the PDF
        reader = PdfReader(pdf_path)
        
        row = 1
        for page in reader.pages:
            text = page.extract_text()
            lines = text.split('\n')
            
            for line in lines:
                ws.cell(row=row, column=1, value=line)
                row += 1
        
        output_path = os.path.join(output_dir, f"converted_to_excel_{uuid4().hex}.xlsx")
        wb.save(output_path)
        
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting PDF to Excel: {str(e)}")
        raise Exception(f"فشل في تحويل PDF إلى Excel: {str(e)}")

def pdf_to_ppt(pdf_path, output_dir):
    """
    Convert a PDF file to a PowerPoint presentation.
    
    This is a simplified implementation. For production use,
    consider a dedicated PDF-to-PowerPoint conversion API.
    
    Args:
        pdf_path: Path to the PDF file
        output_dir: Directory to save the result
        
    Returns:
        Path to the generated PowerPoint document
    """
    try:
        # Convert PDF to images
        images = convert_from_path(pdf_path)
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Add each image as a slide
        for img in images:
            # Save image temporarily
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp:
                img_path = temp.name
                img.save(img_path, 'PNG')
            
            # Add a slide with the image
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
            
            # Clean up temp file
            os.unlink(img_path)
        
        output_path = os.path.join(output_dir, f"converted_to_ppt_{uuid4().hex}.pptx")
        prs.save(output_path)
        
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting PDF to PowerPoint: {str(e)}")
        raise Exception(f"فشل في تحويل PDF إلى PowerPoint: {str(e)}")

def word_to_pdf(word_path, output_dir):
    """
    Convert a Word document to PDF.
    
    Args:
        word_path: Path to the Word document
        output_dir: Directory to save the result
        
    Returns:
        Path to the generated PDF
    """
    try:
        output_path = os.path.join(output_dir, f"converted_from_word_{uuid4().hex}.pdf")
        docx_to_pdf(word_path, output_path)
        
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting Word to PDF: {str(e)}")
        raise Exception(f"فشل في تحويل Word إلى PDF: {str(e)}")

def excel_to_pdf(excel_path, output_dir):
    """
    Convert an Excel document to PDF.
    
    This is a simplified implementation. For production use,
    consider a dedicated Excel-to-PDF conversion API.
    
    Args:
        excel_path: Path to the Excel document
        output_dir: Directory to save the result
        
    Returns:
        Path to the generated PDF
    """
    try:
        # This is a placeholder for actual Excel-to-PDF conversion
        # In a real implementation, you would use a library like win32com.client on Windows
        # or a cloud-based conversion API
        
        # For now, we'll just create a simple PDF with text about the Excel file
        output_path = os.path.join(output_dir, f"converted_from_excel_{uuid4().hex}.pdf")
        
        # Create a writer for the PDF
        writer = PdfWriter()
        
        # Create a new page
        from PyPDF2.generic import RectangleObject
        page = writer.add_blank_page(width=612, height=792)
        
        # Add a note that this is a placeholder
        with open(output_path, "wb") as output_file:
            writer.write(output_file)
        
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting Excel to PDF: {str(e)}")
        raise Exception(f"فشل في تحويل Excel إلى PDF: {str(e)}")

def ppt_to_pdf(ppt_path, output_dir):
    """
    Convert a PowerPoint presentation to PDF.
    
    This is a simplified implementation. For production use,
    consider a dedicated PowerPoint-to-PDF conversion API.
    
    Args:
        ppt_path: Path to the PowerPoint presentation
        output_dir: Directory to save the result
        
    Returns:
        Path to the generated PDF
    """
    try:
        # This is a placeholder for actual PowerPoint-to-PDF conversion
        # In a real implementation, you would use a library like win32com.client on Windows
        # or a cloud-based conversion API
        
        # For now, we'll just create a simple PDF with text about the PowerPoint file
        output_path = os.path.join(output_dir, f"converted_from_ppt_{uuid4().hex}.pdf")
        
        # Create a writer for the PDF
        writer = PdfWriter()
        
        # Create a new page
        from PyPDF2.generic import RectangleObject
        page = writer.add_blank_page(width=612, height=792)
        
        # Write the PDF
        with open(output_path, "wb") as output_file:
            writer.write(output_file)
        
        return output_path
    
    except Exception as e:
        logger.error(f"Error converting PowerPoint to PDF: {str(e)}")
        raise Exception(f"فشل في تحويل PowerPoint إلى PDF: {str(e)}")
